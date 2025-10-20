import os, base64, json, uuid, io, tempfile, shutil, urllib.parse
from typing import Optional, List
from flask import Blueprint, jsonify, request, send_file, send_from_directory, current_app
from sqlalchemy import text
from routes_auth import jwt_required, get_current_user
import boto3
from cryptography.hazmat.primitives.ciphers.aead import AESGCM
from app import db
from utils.file_parser import parse_any
from utils.masker import (
    process_pii, apply_mask_str, ALWAYS_MASK,
    handle_masking, VALID_KEYS, _normalize_allowed_types
)
from utils.crypto_core import encrypt_bytes, decrypt_bytes
from utils.audit_logger import log_audit
from utils.encryptor import encrypt_file, parse_container  # ⬅️ 컨테이너 파서 임포트
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

# ──────────────────────────────
# 기본 설정
# ──────────────────────────────
REGION = os.getenv("AWS_REGION", "ap-northeast-2")
KMS_KEY_ID = os.getenv("KMS_KEY_ID", "alias/KEK")
crypto_bp = Blueprint("crypto", __name__, url_prefix="/api/crypto")
masked_bp = Blueprint("masked_files", __name__)
kms = boto3.client("kms", region_name=REGION)

# ──────────────────────────────
# 유틸
# ──────────────────────────────
def _is_uuid(s: str) -> bool:
    try:
        uuid.UUID(str(s))
        return True
    except Exception:
        return False

MIMETYPES = {
    ".txt": "text/plain",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".csv": "text/csv",
    ".json": "application/json",
    ".rtf": "application/rtf",
    ".odt": "application/vnd.oasis.opendocument.text",
    ".ppt": "application/vnd.ms-powerpoint",
    ".doc": "application/msword",
    ".lkm": "application/octet-stream",
    ".enc": "application/octet-stream",
}
def guess_mime_by_ext(fname: str) -> str:
    ext = os.path.splitext(fname)[1].lower()
    return MIMETYPES.get(ext, "application/octet-stream")

# ──────────────────────────────
# DOCX / PPTX / XLSX 마스킹
# ──────────────────────────────
def mask_docx(in_path: str, out_path: str, allowed_types: Optional[List[str]]):
    doc = DocxDocument(in_path)
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            run.text = apply_mask_str(run.text, allowed_types)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.text = apply_mask_str(run.text, allowed_types)
    doc.save(out_path)

def mask_pptx(in_path: str, out_path: str, allowed_types: Optional[List[str]]):
    prs = Presentation(in_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para.text = apply_mask_str(para.text, allowed_types)
    prs.save(out_path)

def mask_xlsx(in_path: str, out_path: str, allowed_types: Optional[List[str]]):
    wb = openpyxl.load_workbook(in_path, data_only=False, keep_links=False)
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell.value = apply_mask_str(cell.value, allowed_types)
    tmp_out = out_path + ".tmp"
    wb.save(tmp_out)
    wb.close()
    shutil.move(tmp_out, out_path)

# ──────────────────────────────
# 암호화 + 마스킹
# ──────────────────────────────
@crypto_bp.post("/encrypt")
@jwt_required
def encrypt():
    try:
        current_app.logger.info("[encrypt] 시작")

        if "file" not in request.files:
            return jsonify({"ok": False, "error": "file_required"}), 400

        user = get_current_user()
        f = request.files["file"]
        raw = f.read()
        filename = f.filename or "upload.bin"
        name, ext = os.path.splitext(filename)
        token_id = str(uuid.uuid4())

        # 마스킹 타입 처리
        mask_types = None
        if request.form.get("mask_types"):
            try:
                mask_types = json.loads(request.form["mask_types"])
            except Exception:
                mask_types = request.form.getlist("mask_types")

        if not mask_types:
            types_to_use = None
        else:
            norm = _normalize_allowed_types(mask_types)
            types_to_use = sorted(ALWAYS_MASK | set(norm)) or None

        parsed_text = parse_any(filename, raw)
        masked_text = None
        stats = {}
        if parsed_text:
            masked_text, _, stats = process_pii(parsed_text, types_to_use)

        tmp_dir = tempfile.mkdtemp(prefix="locku_")
        try:
            upload_tmp = os.path.join(tmp_dir, f"src{ext or '.bin'}")
            with open(upload_tmp, "wb") as w:
                w.write(raw)

            masked_dir = os.path.join(current_app.root_path, "static", "masked")
            os.makedirs(masked_dir, exist_ok=True)
            masked_name = f"{name}.masked{ext or '.bin'}"
            masked_path = os.path.join(masked_dir, masked_name)

            ext_lower = ext.lower()
            if ext_lower == ".docx":
                mask_docx(upload_tmp, masked_path, types_to_use)
            elif ext_lower == ".pptx":
                mask_pptx(upload_tmp, masked_path, types_to_use)
            elif ext_lower == ".xlsx":
                mask_xlsx(upload_tmp, masked_path, types_to_use)
            elif ext_lower in [".txt", ".csv", ".json", ".rtf", ".odt"]:
                with open(masked_path, "w", encoding="utf-8") as w:
                    w.write(masked_text or parsed_text or "")
            else:
                handle_masking(upload_tmp, masked_path, types_to_use)

            meta = {
                "filename": filename,
                "pii_count": sum(stats.values()) if stats else 0,
                "pii_stats": stats,
                "masked_file": f"/static/masked/{masked_name}",
            }

            enc = encrypt_bytes(raw, meta)  # base64 문자열 반환
            params = {
                "token": token_id,
                "wdek": enc["wdek"],
                "iv": enc["iv"],
                "ciphertext": enc["ciphertext"],
                "tag": (enc.get("tag") or None),
                "meta": json.dumps(meta),
                "user_id": user.id,
            }
            db.session.execute(text("""
                INSERT INTO crypto_store.crypto_tokens
                    (token, wdek, iv, ciphertext, tag, meta, user_id)
                VALUES
                    (:token, :wdek, :iv, :ciphertext, :tag, CAST(:meta AS JSONB), :user_id)
            """), params)
            db.session.commit()

            log_audit("crypto_store", "ENCRYPT", {"status": "SUCCESS", "token_id": token_id, "meta": meta, "user_id": user.id})

            host = request.host or "lockument.duckdns.org"
            scheme = "https" if request.is_secure or host.endswith(":443") else "https"
            base_url = f"{scheme}://{host}"
            masked_url = f"{base_url}/static/masked/{urllib.parse.quote(masked_name)}"

            return jsonify({
                "ok": True,
                "token_id": token_id,
                "token": token_id,  # 하위호환
                "masked_file_url": masked_url,
                "pii_stats": stats,
            }), 200

        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    except Exception as e:
        db.session.rollback()
        current_app.logger.exception(f"[encrypt] 오류: {e}")
        log_audit("crypto_store", "ENCRYPT", {"status": "FAILED", "error": str(e)})
        return jsonify({"ok": False, "error": "encrypt_failed", "message": str(e)}), 500

# ──────────────────────────────
# ✅ 전체 파일 암호화 (KMS AES-GCM) - .enc + .lkm 모두 생성/반환
# ──────────────────────────────
@crypto_bp.post("/encrypt/full")
@jwt_required
def encrypt_full():
    tmp_dir = None
    try:
        current_app.logger.info("[encrypt_full] 시작")

        if "file" not in request.files:
            return jsonify({"ok": False, "error": "file_required"}), 400

        user = get_current_user()
        f = request.files["file"]
        filename = f.filename or "upload.bin"

        tmp_dir = tempfile.mkdtemp(prefix="enc_full_")
        tmp_path = os.path.join(tmp_dir, filename)
        f.save(tmp_path)

        kms_key_id = os.getenv("KMS_KEY_ID")
        result = encrypt_file(tmp_path, user.id, kms_key_id)

        enc_dir = os.path.join(current_app.root_path, "static", "encrypted")
        os.makedirs(enc_dir, exist_ok=True)

        # .enc 파일 저장
        final_enc_path = shutil.move(
            result["enc_path"],
            os.path.join(enc_dir, os.path.basename(result["enc_path"]))
        )

        # .lkm 파일 저장 (컨테이너)
        final_lkm_path = None
        if result.get("lkm_path"):
            final_lkm_path = shutil.move(
                result["lkm_path"],
                os.path.join(enc_dir, os.path.basename(result["lkm_path"]))
            )

        # EncDEK 소유자 검증용 보관
        enc_dek_bytes = base64.b64decode(result["metadata"]["enc_dek"])
        db.session.execute(text("""
            INSERT INTO public.crypto_store
                (kms_key_arn, data_key_ciphertext, created_by, user_id, created_at)
            VALUES
                (:kms_key_arn, :data_key_ciphertext, :created_by, :user_id, CURRENT_TIMESTAMP)
        """), {
            "kms_key_arn": kms_key_id,
            "data_key_ciphertext": enc_dek_bytes,
            "created_by": user.id,
            "user_id": user.id
        })
        db.session.commit()

        enc_download_url = f"/api/crypto/download/{urllib.parse.quote(os.path.basename(final_enc_path))}"
        lkm_download_url = None
        if final_lkm_path:
            lkm_download_url = f"/api/crypto/download/{urllib.parse.quote(os.path.basename(final_lkm_path))}"

        log_audit(
            "crypto_store",
            "ENCRYPT_FILE_FULL",
            {
                "status": "SUCCESS",
                "user_id": user.id,
                "filename": filename,
                "metadata": result["metadata"],
            },
        )

        return jsonify({
            "ok": True,
            "message": "파일이 성공적으로 암호화되었습니다.",
            "encrypted_file": os.path.basename(final_enc_path),  # 예: report.pdf.enc
            "download_url": enc_download_url,
            "lkm_file": os.path.basename(final_lkm_path) if final_lkm_path else None,  # 예: report.pdf.lkm
            "lkm_download_url": lkm_download_url,
            "metadata": result["metadata"]
        }), 200

    except Exception as e:
        db.session.rollback()
        current_app.logger.exception(f"[encrypt_full] 오류: {e}")
        log_audit(
            "crypto_store",
            "ENCRYPT_FILE_FULL",
            {
                "status": "FAILED",
                "user_id": getattr(locals().get("user", None), "id", None),
                "filename": locals().get("filename"),
                "error": str(e),
            },
        )
        return jsonify({"ok": False, "error": "encrypt_full_failed", "message": str(e)}), 500

    finally:
        if tmp_dir:
            shutil.rmtree(tmp_dir, ignore_errors=True)

# ──────────────────────────────
# ✅ 암호화된 파일 다운로드
# ──────────────────────────────
@crypto_bp.get("/download/<path:filename>")
@jwt_required
def download_encrypted_file(filename):
    try:
        decoded = urllib.parse.unquote(filename)
        safe_name = os.path.normpath(decoded).replace("\\", "/")
        if safe_name.startswith("../"):
            return jsonify({"ok": False, "error": "invalid_path"}), 400

        enc_dir = os.path.join(current_app.root_path, "static", "encrypted")
        enc_path = os.path.join(enc_dir, safe_name)
        if not os.path.exists(enc_path):
            return jsonify({"ok": False, "error": "file_not_found"}), 404

        mime = guess_mime_by_ext(safe_name)
        user = get_current_user()
        log_audit(
            "crypto_store",
            "DOWNLOAD_ENCRYPTED_FILE",
            {
                "status": "SUCCESS",
                "user_id": user.id,
                "filename": safe_name,
            },
        )

        return send_from_directory(
            enc_dir,
            safe_name,
            as_attachment=True,
            mimetype=mime,
            download_name=os.path.basename(safe_name)  # 예: report.pdf.enc / report.pdf.lkm
        )

    except Exception as e:
        current_app.logger.exception(f"[download_encrypted_file] 오류: {e}")
        return jsonify({"ok": False, "error": "download_failed", "message": str(e)}), 500

# ──────────────────────────────
# 복호화 (DB 기반 + 소유자 검증)
# ──────────────────────────────
@crypto_bp.post("/decrypt")
@jwt_required
def decrypt():
    try:
        user = get_current_user()
        body = request.get_json(silent=True) or request.form.to_dict() or {}

        token_id = body.get("token_id") or body.get("token")
        if not token_id:
            return jsonify({"ok": False, "error": "token_id_required"}), 400

        if isinstance(token_id, str) and token_id.strip().startswith("{"):
            return jsonify({
                "ok": False,
                "error": "invalid_token_id",
                "message": "token_id must be a UUID (not a JSON bundle)"
            }), 400
        if not _is_uuid(token_id):
            return jsonify({"ok": False, "error": "invalid_token_id_format"}), 400

        row = db.session.execute(text("""
            SELECT user_id, wdek, iv, ciphertext, meta
            FROM crypto_store.crypto_tokens
            WHERE token = :token
        """), {"token": str(token_id)}).fetchone()

        if not row:
            return jsonify({"ok": False, "error": "token_not_found"}), 404

        owner_id, wdek_b64, iv_b64, ct_b64, raw_meta = row
        if owner_id != user.id:
            log_audit("crypto_store", "DECRYPT_DENIED", {
                "status": "FAILED",
                "reason": "unauthorized_user",
                "token_id": str(token_id),
                "owner_id": owner_id,
                "user_id": user.id
            })
            return jsonify({"ok": False, "error": "unauthorized_user"}), 403

        meta = json.loads(raw_meta) if isinstance(raw_meta, str) else (raw_meta or {})
        # ✅ 원본 파일명 그대로 내려줌
        filename = meta.get("filename", "decrypted_file.bin")
        mime = guess_mime_by_ext(filename)

        plaintext = decrypt_bytes(wdek_b64, iv_b64, ct_b64)
        log_audit("crypto_store", "DECRYPT", {"status": "SUCCESS", "token_id": str(token_id), "meta": meta, "user_id": user.id})

        bio = io.BytesIO(plaintext)
        bio.seek(0)
        response = send_file(bio, mimetype=mime, as_attachment=True, download_name=filename)
        response.headers["Content-Disposition"] = f"attachment; filename*=UTF-8''{urllib.parse.quote(filename)}"
        return response

    except Exception as e:
        db.session.rollback()
        current_app.logger.exception("[decrypt] 실패")
        log_audit("crypto_store", "DECRYPT", {"status": "FAILED", "error": str(e)})
        return jsonify({"ok": False, "error": "decrypt_failed", "message": str(e)}), 500

# ──────────────────────────────
# ✅ 전체 복호화 (KMS AES-GCM + 소유자 검증)
#    - metadata 없으면 .lkm 컨테이너에서 헤더 파싱하여 복호화
#    - 복호화 결과는 원래 파일명/확장자 그대로 다운로드
# ──────────────────────────────
@crypto_bp.post("/decrypt/full")
@jwt_required
def decrypt_full():
    tmp_dir = None
    try:
        current_app.logger.info("[decrypt_full] 시작")

        user = get_current_user()
        file = request.files.get("file")
        meta_str = request.form.get("metadata")  # 과거 방식 호환

        if not file:
            return jsonify({"ok": False, "error": "file_required"}), 400

        tmp_dir = tempfile.mkdtemp(prefix="dec_full_")
        enc_path = os.path.join(tmp_dir, file.filename)
        file.save(enc_path)

        with open(enc_path, "rb") as rf:
            blob = rf.read()

        plaintext = None
        original_name = None  # ✅ 복호화 결과 파일명

        # 1) metadata가 없으면: 컨테이너(.lkm) 파싱 시도
        if not meta_str:
            parsed = parse_container(blob)  # -> (header(dict), ciphertext:bytes) or None
            if parsed:
                metadata, ciphertext = parsed

                # 파일명 헤더가 있으면 사용, 없으면 업로드 파일명에서 .lkm 제거
                original_name = metadata.get("filename")
                if not original_name:
                    original_name = file.filename[:-4] if file.filename.lower().endswith(".lkm") else file.filename

                # 소유자 검증: EncDEK로 DB 대조
                enc_dek_bytes = base64.b64decode(metadata["enc_dek"])
                row = db.session.execute(text("""
                    SELECT user_id FROM public.crypto_store
                    WHERE data_key_ciphertext = :data_key_ciphertext
                """), {"data_key_ciphertext": enc_dek_bytes}).fetchone()
                if row and row[0] != user.id:
                    log_audit(
                        "crypto_store",
                        "DECRYPT_FILE_FULL_DENIED",
                        {"status": "FAILED", "owner_id": row[0], "user_id": user.id},
                    )
                    return jsonify({"ok": False, "error": "unauthorized_user"}), 403

                # KMS 복호화 → AES-GCM 복호화
                resp = kms.decrypt(CiphertextBlob=enc_dek_bytes)
                plain_key = resp["Plaintext"]
                iv = base64.b64decode(metadata["iv"])
                aesgcm = AESGCM(plain_key)
                plaintext = aesgcm.decrypt(iv, ciphertext, None)

            else:
                # 컨테이너가 아니면 과거 방식 요구
                return jsonify({"ok": False, "error": "missing_metadata", "message": "컨테이너 파일이 아니면 metadata가 필요합니다."}), 400

        # 2) metadata가 있으면: 기존(.enc + metadata) 방식
        else:
            metadata = json.loads(meta_str)
            enc_dek = base64.b64decode(metadata["enc_dek"])
            iv = base64.b64decode(metadata["iv"])

            # 업로드 파일명에서 .enc 제거해 원본 파일명 복원
            original_name = file.filename[:-4] if file.filename.lower().endswith(".enc") else file.filename

            row = db.session.execute(text("""
                SELECT user_id FROM public.crypto_store
                WHERE data_key_ciphertext = :data_key_ciphertext
            """), {"data_key_ciphertext": enc_dek}).fetchone()
            if row and row[0] != user.id:
                log_audit(
                    "crypto_store",
                    "DECRYPT_FILE_FULL_DENIED",
                    {"status": "FAILED", "owner_id": row[0], "user_id": user.id},
                )
                return jsonify({"ok": False, "error": "unauthorized_user"}), 403

            resp = kms.decrypt(CiphertextBlob=enc_dek)
            plain_key = resp["Plaintext"]
            aesgcm = AESGCM(plain_key)
            ciphertext = blob
            plaintext = aesgcm.decrypt(iv, ciphertext, None)

        # ✅ 원래 파일명/확장자 그대로 반환
        dec_name = original_name or "decrypted_file.bin"
        mime = guess_mime_by_ext(dec_name)

        dec_path = os.path.join(tmp_dir, dec_name)
        with open(dec_path, "wb") as wf:
            wf.write(plaintext)

        log_audit(
            "crypto_store",
            "DECRYPT_FILE_FULL",
            {"status": "SUCCESS", "user_id": user.id, "filename": file.filename},
        )

        return send_file(dec_path, as_attachment=True, download_name=dec_name, mimetype=mime)

    except Exception as e:
        current_app.logger.exception(f"[decrypt_full] 오류: {e}")
        log_audit(
            "crypto_store",
            "DECRYPT_FILE_FULL",
            {
                "status": "FAILED",
                "user_id": getattr(locals().get("user", None), "id", None),
                "filename": getattr(locals().get("file", None), "filename", None),
                "error": str(e),
            },
        )
        return jsonify({"ok": False, "error": "decrypt_full_failed", "message": str(e)}), 500

    finally:
        if tmp_dir:
            shutil.rmtree(tmp_dir, ignore_errors=True)

# ──────────────────────────────
# 감사 로그
# ──────────────────────────────
@crypto_bp.get("/audit/recent")
def recent_audits():
    try:
        ip = request.headers.get("X-Forwarded-For", request.remote_addr)
        ua = request.headers.get("User-Agent", "unknown")
        user = None
        try:
            user = get_current_user()
        except Exception:
            pass

        sql = text("""
            SELECT op, row_new->>'status' AS status, changed_at
            FROM pii_audit.audit_logs
            ORDER BY changed_at DESC
            LIMIT 20;
        """)
        rows = db.session.execute(sql).mappings().all()

        logs = []
        for r in rows:
            entry = dict(r)
            entry["ip"] = ip
            entry["user_agent"] = ua
            entry["user"] = getattr(user, "username", None)
            logs.append(entry)

        return jsonify({"ok": True, "logs": logs})
    except Exception as e:
        current_app.logger.exception(f"[audit/recent] failed: {e}")
        return jsonify({"ok": False, "error": str(e)}), 500

# ──────────────────────────────
# 최근 암호화 파일 목록
# ──────────────────────────────
@crypto_bp.get("/recent")
def list_recent_files():
    limit = request.args.get("limit", 20, type=int)
    sql = text("""
        SELECT meta->>'filename' AS filename, created_at
        FROM crypto_store.crypto_tokens
        ORDER BY created_at DESC
        LIMIT :limit;
    """)
    rows = db.session.execute(sql, {"limit": limit}).mappings().all()
    files = [
        {"filename": r["filename"], "created_at": r["created_at"].isoformat() if r["created_at"] else None}
        for r in rows
    ]
    return jsonify({"ok": True, "files": files})

# ──────────────────────────────
# ✅ 마스킹 파일 다운로드 (퍼블릭 경로)
# ──────────────────────────────
@masked_bp.route("/static/masked/<path:filename>")
def serve_masked_public(filename):
    decoded = urllib.parse.unquote(filename)
    safe = os.path.normpath(decoded).replace("\\", "/")
    if safe.startswith("../"):
        return jsonify({"ok": False, "error": "invalid_path"}), 400

    masked_dir = os.path.join(current_app.root_path, "static", "masked")
    file_path = os.path.join(masked_dir, safe)
    if not os.path.exists(file_path):
        return jsonify({"ok": False, "error": "file_not_found"}), 404

    mime = guess_mime_by_ext(safe)
    return send_from_directory(
        masked_dir,
        safe,
        as_attachment=True,
        mimetype=mime,
        download_name=os.path.basename(safe)
    )

