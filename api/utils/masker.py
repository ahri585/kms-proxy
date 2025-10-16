import os, re, json, hashlib, openpyxl, pdfplumber
from pptx import Presentation
from docx import Document as DocxDocument
from typing import Optional, List, Tuple

def sha256_hex(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()

AUTH_VALUE_SUB_RE = re.compile(
    r'((password|api[_ ]?key|token)\s*[:=]\s*)([\'"]?)[^\'",\s]+([\'"]?)',
    re.IGNORECASE
)

PII_RULES_MAP = {
    "rrn": (
        re.compile(r"\b\d{6}[-]?\d{7}\b"),
        lambda s: "######-*******"
    ),
    "email": (
        re.compile(r"\b([A-Za-z0-9._%+-])([A-Za-z0-9._%+-]*)(@[A-Za-z0-9.-]+\.[A-Za-z]{2,})\b"),
        lambda s: (s[0] + "***" + s[s.find('@'):]) if "@" in s else s
    ),
    "card_or_acct": (
        re.compile(r"\b\d{10,19}\b"),
        lambda s: (s[:6] + "******" + s[-4:]) if len(s) >= 10 else "******"
    ),
    "auth": (
        re.compile(r"\b(password|api[_ ]?key|token)\s*[:=]\s*['\"]?[^'\",\s]+['\"]?", re.IGNORECASE),
        lambda s: AUTH_VALUE_SUB_RE.sub(r"\1[SECRET]", s)
    ),
    "passport": (
        re.compile(r"\b[A-Z][0-9]{8}\b", re.IGNORECASE),
        lambda s: s[0] + "********"
    ),
    "license": (
        re.compile(r"\b\d{2}-\d{2}-\d{6}-\d{2}\b|\b\d{2}-\d{6}-\d{2}-\d{2}\b"),
        lambda s: re.sub(r"\d", "*", s)
    ),
    "foreign_id": (
        re.compile(r"\b\d{6}-[5-8]\d{6}\b"),
        lambda s: "######-*******"
    ),
    "phone": (
        re.compile(r"\b(01[016789])[ ._-]?\d{3,4}[ ._-]?\d{4}\b"),
        lambda s: "010-****-****"
    ),
    "address": (
        re.compile(r"[가-힣]+\s*(시|군|구)\s*[가-힣0-9\s\-]*(동|읍|면|리|로|길)\s*\d*[-]?\d*호?"),
        lambda s: "[주소마스킹]"
    ),
    "name": (
        re.compile(r"(?<![가-힣])([가-힣]{2,4})(?![가-힣])"),
        lambda s: "[이름마스킹]"
    ),
}

ALWAYS_MASK = {"rrn", "passport", "license", "foreign_id", "auth"}
VALID_KEYS = set(PII_RULES_MAP.keys())

ORDERED_LABELS = [
    "auth",
    "email",
    "phone",
    "address",
    "rrn",
    "foreign_id",
    "passport",
    "license",
    "card_or_acct",
    "name",
]

def process_pii(text: str, allowed_types: Optional[List[str]] = None) -> Tuple[str, list, dict]:
    hits, stats, masked = [], {}, text
    sel = _normalize_allowed_types(allowed_types)
    for a in ALWAYS_MASK:
        if a not in sel:
            sel.append(a)

    types_to_use = [lbl for lbl in ORDERED_LABELS if lbl in sel]

    for label in types_to_use:
        pattern, mask_fn = PII_RULES_MAP[label]

        def repl(m):
            orig = m.group(0)
            masked_val = mask_fn(orig)
            hits.append({
                "type": label,
                "masked": masked_val,
                "sha256": sha256_hex(orig.encode("utf-8", "ignore")),
            })
            stats[label] = stats.get(label, 0) + 1
            return masked_val

        masked = pattern.sub(repl, masked)
    return masked, hits, stats

def apply_mask_str(s: str, allowed_types: Optional[List[str]]) -> str:
    return process_pii(s, allowed_types)[0]

def handle_masking(src_path: str, dst_masked_path: str, allowed_types: Optional[List[str]]) -> str:
    ext = os.path.splitext(src_path)[1].lower()
    os.makedirs(os.path.dirname(dst_masked_path), exist_ok=True)

    if ext in [".txt", ".csv", ".json"]:
        with open(src_path, "r", encoding="utf-8", errors="ignore") as f:
            masked = apply_mask_str(f.read(), allowed_types)
        with open(dst_masked_path, "w", encoding="utf-8") as f:
            f.write(masked)

    elif ext == ".xlsx":
        wb = openpyxl.load_workbook(src_path, data_only=True)
        out = []
        for s in wb.sheetnames:
            for row in wb[s].iter_rows(values_only=True):
                out.append(" ".join(str(c) for c in row if c is not None))
        masked_txt = apply_mask_str("\n".join(out), allowed_types)
        base, _ = os.path.splitext(dst_masked_path)
        dst_masked_path = f"{base}_masked.txt"
        with open(dst_masked_path, "w", encoding="utf-8") as f:
            f.write(masked_txt)

    elif ext == ".pdf":
        text = []
        with pdfplumber.open(src_path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text() or "")
        masked_txt = apply_mask_str("\n".join(text), allowed_types)
        base, _ = os.path.splitext(dst_masked_path)
        dst_masked_path = f"{base}_masked.txt"
        with open(dst_masked_path, "w", encoding="utf-8") as f:
            f.write(masked_txt)

    elif ext == ".docx":
        doc = DocxDocument(src_path)
        for p in doc.paragraphs:
            p.text = apply_mask_str(p.text, allowed_types)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        p.text = apply_mask_str(p.text, allowed_types)
        doc.save(dst_masked_path)

    elif ext == ".pptx":
        prs = Presentation(src_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para.text = apply_mask_str(para.text, allowed_types)
        prs.save(dst_masked_path)

    else:
        with open(src_path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        masked = apply_mask_str(raw, allowed_types)
        base, _ = os.path.splitext(dst_masked_path)
        dst_masked_path = f"{base}_masked.txt"
        with open(dst_masked_path, "w", encoding="utf-8") as f:
            f.write(masked)

    return os.path.basename(dst_masked_path)

ALIASES = {
    "전화": "phone",
    "이메일": "email",
    "카드": "card_or_acct",
    "계좌": "card_or_acct",
    "비밀번호": "auth",
    "주민등록번호": "rrn",
    "여권": "passport",
    "운전면허": "license",
    "외국인등록": "foreign_id",
    "주소": "address",
    "이름": "name",
}

def _normalize_allowed_types(allowed_types: Optional[List[str] | str]) -> List[str]:
    if not allowed_types:
        return list(ALWAYS_MASK)

    if isinstance(allowed_types, str):
        parts = [p.strip() for p in allowed_types.split(",") if p.strip()]
    else:
        parts = []
        for x in allowed_types:
            if x is None:
                continue
            s = str(x)
            if "," in s:
                parts.extend([p.strip() for p in s.split(",") if p.strip()])
            else:
                parts.append(s.strip())

    out = []
    for p in parts:
        k = ALIASES.get(p.lower(), p.lower())
        if k in VALID_KEYS:
            out.append(k)
    return list(dict.fromkeys(out))
