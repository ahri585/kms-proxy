import io, json, pdfplumber, docx, openpyxl, pandas as pd
from pptx import Presentation

def parse_any(filename: str, b: bytes) -> str:
    """파일에서 텍스트 추출 (PDF, DOCX, PPTX, XLSX 등)"""
    fn = filename.lower()
    try:
        if fn.endswith(".pdf"):
            with pdfplumber.open(io.BytesIO(b)) as pdf:
                return "\n".join((p.extract_text() or "") for p in pdf.pages)
        if fn.endswith(".docx"):
            d = docx.Document(io.BytesIO(b))
            return "\n".join(p.text for p in d.paragraphs if p.text)
        if fn.endswith(".xlsx"):
            wb = openpyxl.load_workbook(io.BytesIO(b), data_only=True)
            out = []
            for s in wb.sheetnames:
                for row in wb[s].iter_rows(values_only=True):
                    out.append(" ".join(str(c) for c in row if c is not None))
            return "\n".join(out)
        if fn.endswith(".pptx"):
            prs = Presentation(io.BytesIO(b))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text:
                                out.append(para.text)
            return "\n".join(out)
        if fn.endswith(".txt"):
            return b.decode("utf-8", "ignore")
        if fn.endswith(".csv"):
            try:
                return pd.read_csv(io.BytesIO(b), encoding="utf-8").to_string(index=False)
            except Exception:
                return pd.read_csv(io.BytesIO(b), encoding="cp949").to_string(index=False)
        if fn.endswith(".json"):
            return json.dumps(json.loads(b.decode("utf-8", "ignore")), ensure_ascii=False, indent=2)
    except Exception:
        return b.decode("utf-8", "ignore")
    return ""

